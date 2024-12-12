import argparse
from azure.ai.translation.text import TextTranslationClient
from azure.core.credentials import AzureKeyCredential
from pptx import Presentation

def translate_presentation(presentation, client, target_language_code):
    for slide_index, slide in enumerate(presentation.slides, start=1):
        print(f"Translating Slide {slide_index} of {len(presentation.slides)}...")
        
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        translate_text_frame(cell.text_frame, client, target_language_code)
            elif shape.has_text_frame:
                translate_text_frame(shape.text_frame, client, target_language_code)
        
        if slide.has_notes_slide:
            translate_text_frame(slide.notes_slide.notes_text_frame, client, target_language_code)

def translate_text_frame(text_frame, client, target_language_code):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                try:
                    body = [{"text": run.text}]
                    response = client.translate(body=body, to_language=[target_language_code])
                    translated_text = response[0]["translations"][0]["text"]
                    run.text = translated_text
                except Exception as e:
                    print(f"Translation failed for text '{run.text}': {e}")

def main():
    argument_parser = argparse.ArgumentParser(
        description="Translate PPTX files using Azure Translator API"
    )
    argument_parser.add_argument(
        "source_language_code", type=str,
        help="The language code for the language of the source text. Example: en"
    )
    argument_parser.add_argument(
        "target_language_code", type=str,
        help="The language code for the target text. Example: es"
    )
    argument_parser.add_argument(
        "input_file_path", type=str,
        help="The path to the PPTX file to translate"
    )
    argument_parser.add_argument(
        "--azure_key", type=str, required=True,
        help="Your Azure Translator API key"
    )
    argument_parser.add_argument(
        "--endpoint", type=str, required=True,
        help="Your Azure Translator API endpoint"
    )
    args = argument_parser.parse_args()

    client = TextTranslationClient(endpoint=args.endpoint, credential=AzureKeyCredential(args.azure_key))

    print(f"Translating {args.input_file_path} from {args.source_language_code} to {args.target_language_code}...")
    presentation = Presentation(args.input_file_path)

    translate_presentation(presentation, client, args.target_language_code)

    output_file_path = args.input_file_path.replace(".pptx", f"-{args.target_language_code}.pptx")
    print(f"Saving translated presentation to {output_file_path}...")
    presentation.save(output_file_path)

if __name__ == '__main__':
    main()